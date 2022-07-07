from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from bs4 import BeautifulSoup
import zipfile
import script_zip
import shutil
import re
from PIL import Image

table_style_dict = { 'NoStyleNoGrid' : '{2D5ABB26-0587-4C30-8999-92F81FD0307C}', 'ThemedStyle1Accent1' : '{3C2FFA5D-87B4-456A-9821-1D502468CF0F}', 'ThemedStyle1Accent2' : '{284E427A-3D55-4303-BF80-6455036E1DE7}', 'ThemedStyle1Accent3' : '{69C7853C-536D-4A76-A0AE-DD22124D55A5}', 'ThemedStyle1Accent4' : '{775DCB02-9BB8-47FD-8907-85C794F793BA}', 'ThemedStyle1Accent5' : '{35758FB7-9AC5-4552-8A53-C91805E547FA}', 'ThemedStyle1Accent6' : '{08FB837D-C827-4EFA-A057-4D05807E0F7C}', 'NoStyleTableGrid' : '{5940675A-B579-460E-94D1-54222C63F5DA}', 'ThemedStyle2Accent1' : '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}', 'ThemedStyle2Accent2' : '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}', 'ThemedStyle2Accent3' : '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}', 'ThemedStyle2Accent4' : '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}', 'ThemedStyle2Accent5' : '{327F97BB-C833-4FB7-BDE5-3F7075034690}', 'ThemedStyle2Accent6' : '{638B1855-1B75-4FBE-930C-398BA8C253C6}', 'LightStyle1' : '{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}', 'LightStyle1Accent1' : '{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}', 'LightStyle1Accent2' : '{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}', 'LightStyle1Accent3' : '{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}', 'LightStyle1Accent4' : '{D27102A9-8310-4765-A935-A1911B00CA55}', 'LightStyle1Accent5' : '{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}', 'LightStyle1Accent6' : '{68D230F3-CF80-4859-8CE7-A43EE81993B5}', 'LightStyle2' : '{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}', 'LightStyle2Accent1' : '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}', 'LightStyle2Accent2' : '{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}', 'LightStyle2Accent3' : '{F2DE63D5-997A-4646-A377-4702673A728D}', 'LightStyle2Accent4' : '{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}', 'LightStyle2Accent5' : '{5A111915-BE36-4E01-A7E5-04B1672EAD32}', 'LightStyle2Accent6' : '{912C8C85-51F0-491E-9774-3900AFEF0FD7}', 'LightStyle3' : '{616DA210-FB5B-4158-B5E0-FEB733F419BA}', 'LightStyle3Accent1' : '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}', 'LightStyle3Accent2' : '{5DA37D80-6434-44D0-A028-1B22A696006F}', 'LightStyle3Accent3' : '{8799B23B-EC83-4686-B30A-512413B5E67A}', 'LightStyle3Accent4' : '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}', 'LightStyle3Accent5' : '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}', 'LightStyle3Accent6' : '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}', 'MediumStyle1' : '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}', 'MediumStyle1Accent1' : '{B301B821-A1FF-4177-AEE7-76D212191A09}', 'MediumStyle1Accent2' : '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}', 'MediumStyle1Accent3' : '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}', 'MediumStyle1Accent4' : '{1E171933-4619-4E11-9A3F-F7608DF75F80}', 'MediumStyle1Accent5' : '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}', 'MediumStyle1Accent6' : '{10A1B5D5-9B99-4C35-A422-299274C87663}', 'MediumStyle2' : '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}', 'MediumStyle2Accent1' : '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}', 'MediumStyle2Accent2' : '{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}', 'MediumStyle2Accent3' : '{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}', 'MediumStyle2Accent4' : '{00A15C55-8517-42AA-B614-E9B94910E393}', 'MediumStyle2Accent5' : '{7DF18680-E054-41AD-8BC1-D1AEF772440D}', 'MediumStyle2Accent6' : '{93296810-A885-4BE3-A3E7-6D5BEEA58F35}', 'MediumStyle3' : '{8EC20E35-A176-4012-BC5E-935CFFF8708E}', 'MediumStyle3Accent1' : '{6E25E649-3F16-4E02-A733-19D2CDBF48F0}', 'MediumStyle3Accent2' : '{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}', 'MediumStyle3Accent3' : '{EB344D84-9AFB-497E-A393-DC336BA19D2E}', 'MediumStyle3Accent4' : '{EB9631B5-78F2-41C9-869B-9F39066F8104}', 'MediumStyle3Accent5' : '{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}', 'MediumStyle3Accent6' : '{2A488322-F2BA-4B5B-9748-0D474271808F}', 'MediumStyle4' : '{D7AC3CCA-C797-4891-BE02-D94E43425B78}', 'MediumStyle4Accent1' : '{69CF1AB2-1976-4502-BF36-3FF5EA218861}', 'MediumStyle4Accent2' : '{8A107856-5554-42FB-B03E-39F5DBC370BA}', 'MediumStyle4Accent3' : '{0505E3EF-67EA-436B-97B2-0124C06EBD24}', 'MediumStyle4Accent4' : '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}', 'MediumStyle4Accent5' : '{22838BEF-8BB2-4498-84A7-C5851F593DF1}', 'MediumStyle4Accent6' : '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}', 'DarkStyle1' : '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}', 'DarkStyle1Accent1' : '{125E5076-3810-47DD-B79F-674D7AD40C01}', 'DarkStyle1Accent2' : '{37CE84F3-28C3-443E-9E96-99CF82512B78}', 'DarkStyle1Accent3' : '{D03447BB-5D67-496B-8E87-E561075AD55C}', 'DarkStyle1Accent4' : '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}', 'DarkStyle1Accent5' : '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}', 'DarkStyle1Accent6' : '{AF606853-7671-496A-8E4F-DF71F8EC918B}', 'DarkStyle2' : '{5202B0CA-FC54-4496-8BCA-5EF66A818D29}', 'DarkStyle2Accent1Accent2' : '{0660B408-B3CF-4A94-85FC-2B1E0A45F4A2}', 'DarkStyle2Accent3Accent4' : '{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}', 'DarkStyle2Accent5Accent6' : '{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}'}



def delete_slide_master(prs, idx):
    '''指定した番号のスライドマスターを削除'''
    xml_slides = prs.slide_masters._sldMasterIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[idx])

def delete_slide_layout(prs, idx):
    '''指定した番号のスライドマスターを削除'''
    xml_layouts= prs.slide_layouts._sldLayoutIdLst
    layouts = list(xml_layouts)
    xml_layouts.remove(layouts[idx])

def get_soup(filepath, target_filepath):
  # xml解凍
  with zipfile.ZipFile(filepath) as zf:
      zf.extract(target_filepath, '')

  # xml読み込み
  with open(target_filepath, 'r') as f:
    soup = BeautifulSoup(f, 'xml')
  return soup

def delete_layouts(filepath, layout_nums: list, master_num = 1):
  '''スライドマスターのレイアウトを複数削除'''
  for layout_num in layout_nums:
    delete_layout(filepath, layout_num, master_num = master_num)

def delete_layout(filepath, layout_num, master_num = 1):
  '''スライドマスターのレイアウトを削除'''
  target_filepath = 'ppt/slideMasters/_rels/slideMaster' + str(master_num) + '.xml.rels'

  # slideMasterの該当レイアウトに関する項目を削除
  soup = get_soup(filepath, target_filepath)

  Id = 'rId'+str(layout_num)
  item = soup.find('Relationship', {'Id': Id})

  # 削除するレイアウト番号を取得
  Id_List = [i.get_attribute_list('Id')[0] for i in soup.find_all('Relationship')]
  Id_List = {int(re.search(r'\d+', value)[0]): value for value in Id_List}
  Id_List = [Id_List[key] for key in sorted(Id_List.keys())]
  Id_arg = [idx for idx, i in enumerate(Id_List) if i == Id][0]

  # slideMaster1のレイアウト本体を削除
  prs_new = Presentation(filepath)  
  delete_slide_layout(prs_new, Id_arg)
  prs_new.save(filepath)
  
  if not item == None:
    # xml編集/保存
    item.extract()
    with open(target_filepath, 'w') as f:
          f.write(str(soup))

    # pptx内xmlファイルの更新
    script_zip.updateZip(filepath, target_filepath, target_filepath)

    # 解凍・編集したxmlを削除
    shutil.rmtree(target_filepath.split('/')[0])
  else:
    print('not found: slideLayout{}'.format(layout_num))



def replace_xml(filepath, replace_dict):
    # pptx内xmlファイルの更新
    for target_filepath, base_filepath in replace_dict.items():
        script_zip.updateZip(filepath, target_filepath, base_filepath)


def change_color_palette(filepath, color_dict, master_num = 1, theme_num = 1):
    # colorpaletteの変更
    # https://www.cresco.co.jp/blog/entry/740/  
    target_filepath = 'ppt/slideMasters/slideMaster' + str(master_num) + '.xml'
    target_filepath = 'ppt/theme/theme' + str(theme_num) + '.xml'
    soup = get_soup(filepath, target_filepath)
    for key, value in color_dict.items():
        if not value == None:
            soup.find_all(key)[0].find('srgbClr')['val'] = value


    with open(target_filepath, 'w') as f:
        f.write(str(soup))

    # pptx内xmlファイルの更新
    script_zip.updateZip(filepath, target_filepath, target_filepath)

    # 解凍・編集したxmlを削除
    shutil.rmtree(target_filepath.split('/')[0])

def add_img(slide, img_path, left, top, width, height):
    """
    args:
        slide[slide]: Slide object
        img_path[str] : Image file path
        left[int]: Position from the left end
        top[int] : Position from top
        width[int]: Width of object
        height[int]: Height of object 
    return:
        None       
    """
    pic = slide.shapes.add_picture(img_path, 0, 0)
    pic.width = width
    pic.height = height
    pic.left = left
    pic.top = top

def get_img_size(img_path):
    img = Image.open(img_path)
    return img.width, img.height

def add_text(p, text, size = 14, bold = True, level = 0, line_spacing = None, alignment = None):
    '''
    text追加関数

    alignment
    CENTER: Center align
    DISTRIBUTE: Evenly distributes e.g. Japanese characters from left to right within a line
    JUSTIFY: Justified, i.e. each line both begins and ends at the margin with spacing between words adjusted such that the line exactly fills the width of the paragraph.
    JUSTIFY_LOW: Justify using a small amount of space between words.
    LEFT: Left aligned
    RIGHT: Right aligned
    THAI_DISTRIBUTE: Thai distributed
    MIXED
    '''
    p.text = text
    p.level = level  # down the bullet level    
    p.font.size = Pt(size)  # font size
    p.font.bold = bold  # font bold
    if not line_spacing == None:
        p.line_spacing = Pt(line_spacing)

    if alignment == 'CENTER':
        p.alignment = PP_ALIGN.CENTER
    elif alignment == 'DISTRIBUTE':
        p.alignment = PP_ALIGN.DISTRIBUTE
    elif alignment == 'JUSTIFY':
        p.alignment = PP_ALIGN.JUSTIFY
    elif alignment == 'JUSTIFY_LOW':
        p.alignment = PP_ALIGN.JUSTIFY_LOW
    elif alignment == 'LEFT':
        p.alignment = PP_ALIGN.LEFT
    elif alignment == 'RIGHT':
        p.alignment = PP_ALIGN.RIGHT
    elif alignment == 'THAI_DISTRIBUTE':
        p.alignment = PP_ALIGN.THAI_DISTRIBUTE
    elif alignment == 'MIXED':
        p.alignment = PP_ALIGN.MIXED


def add_table(slide, data, left, top, width, height, font_size):
    # tableの行数と列数(tableのサイズ)
    rows = len(data)
    cols = len(data[0])

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 表の各セルの中身を記入
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j]
            tfrm = cell.text_frame
            tfrm.paragraphs[0].font.size = Pt(font_size)
            tfrm.paragraphs[0].font.name = 'メイリオ'
            tfrm.paragraphs[0].alignment = PP_ALIGN.CENTER #左右中央揃え
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # tfrm.paragraphs[0].font.color.rgb = RGBColor(12, 34, 56)
            
    # tableの高さを再調整
    table.rows[0].height = int(height/3)
    table.rows[1].height = int(height/3)
    table.rows[2].height = int(height/3)

    # tableの幅を再調整
    table.columns[0].width = int(width/3)
    table.columns[1].width = int(width/3)
    table.columns[2].width = int(width/3)

    tbl = table._graphic_frame.element.graphic.graphicData.tbl
    tbl[0][-1].text = table_style_dict['LightStyle2']